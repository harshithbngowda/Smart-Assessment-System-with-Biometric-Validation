"""
Universal Handwriting Recognition - Supports Images, PDFs, Word, PowerPoint
Works with Google Gemini Vision API
"""

from __future__ import annotations
import google.generativeai as genai
from PIL import Image
from pathlib import Path
import sys
import os
import io


def convert_pdf_to_images(pdf_path: str) -> list[Image.Image]:
    """Convert PDF pages to images"""
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(pdf_path)
        return images
    except ImportError:
        print("❌ pdf2image not installed!")
        print("   Install: pip install pdf2image")
        print("   Also requires poppler: https://github.com/oschwartz10612/poppler-windows/releases/")
        return []


def convert_docx_to_images(docx_path: str) -> list[Image.Image]:
    """Extract images from Word document"""
    try:
        from docx import Document
        from docx.oxml.text.paragraph import CT_P
        from docx.oxml.shape import CT_Picture
        
        doc = Document(docx_path)
        images = []
        
        # Extract embedded images
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                image = Image.open(io.BytesIO(image_data))
                images.append(image)
        
        return images
    except ImportError:
        print("❌ python-docx not installed!")
        print("   Install: pip install python-docx")
        return []


def convert_pptx_to_images(pptx_path: str) -> list[Image.Image]:
    """Extract images from PowerPoint presentation"""
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        images = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image_data = shape.image.blob
                    image = Image.open(io.BytesIO(image_data))
                    images.append(image)
        
        return images
    except ImportError:
        print("❌ python-pptx not installed!")
        print("   Install: pip install python-pptx")
        return []


def recognize_any_file(
    file_path: str,
    api_key: str = None,
    output_path: str = None
) -> str:
    """
    Recognize handwritten text from any file type.
    
    Supported formats:
    - Images: JPG, PNG, WEBP, GIF, BMP
    - PDF: Converts each page to image
    - Word (DOCX): Extracts embedded images
    - PowerPoint (PPTX): Extracts images from slides
    
    Args:
        file_path: Path to the file
        api_key: Google Gemini API key
        output_path: Optional output file path
    
    Returns:
        Recognized text
    """
    print("="*70)
    print("UNIVERSAL HANDWRITING RECOGNITION - Google Gemini")
    print("="*70)
    
    # Get API key
    if api_key is None:
        api_key = os.environ.get('GEMINI_API_KEY')
    
    if not api_key:
        print("\n❌ No API key found!")
        print("   Get FREE API key: https://makersuite.google.com/app/apikey")
        return ""
    
    # Configure Gemini
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel('models/gemini-2.5-flash')
    
    # Detect file type
    file_path = Path(file_path)
    ext = file_path.suffix.lower()
    
    print(f"\nFile: {file_path}")
    print(f"Type: {ext}")
    
    # Convert to images based on file type
    images = []
    
    if ext in ['.jpg', '.jpeg', '.png', '.webp', '.gif', '.bmp']:
        # Direct image file
        images = [Image.open(file_path)]
        print("✅ Image file detected")
    
    elif ext == '.pdf':
        print("📄 PDF detected - converting pages to images...")
        images = convert_pdf_to_images(str(file_path))
        print(f"✅ Extracted {len(images)} pages")
    
    elif ext == '.docx':
        print("📝 Word document detected - extracting images...")
        images = convert_docx_to_images(str(file_path))
        print(f"✅ Extracted {len(images)} images")
    
    elif ext == '.pptx':
        print("📊 PowerPoint detected - extracting images...")
        images = convert_pptx_to_images(str(file_path))
        print(f"✅ Extracted {len(images)} images")
    
    else:
        print(f"❌ Unsupported file type: {ext}")
        print("   Supported: .jpg, .png, .webp, .pdf, .docx, .pptx")
        return ""
    
    if not images:
        print("❌ No images found to process!")
        return ""
    
    # Process each image
    prompt = """
    Please read all the handwritten text in this image and transcribe it exactly as written.
    Preserve line breaks and formatting. Output ONLY the transcribed text.
    """
    
    all_text = []
    
    for i, image in enumerate(images, 1):
        print(f"\n🤖 Processing page/image {i}/{len(images)}...")
        
        try:
            response = model.generate_content([prompt, image])
            text = response.text.strip()
            all_text.append(f"--- Page {i} ---\n{text}\n")
        except Exception as e:
            print(f"⚠️  Error processing page {i}: {e}")
            all_text.append(f"--- Page {i} ---\n[Error processing this page]\n")
    
    # Combine all text
    full_text = "\n".join(all_text)
    
    # Display result
    print("\n" + "="*70)
    print("RECOGNIZED TEXT:")
    print("="*70)
    print(full_text)
    print("="*70)
    
    # Save output
    if output_path:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(full_text)
        
        print(f"\n✅ Saved to: {output_path}")
    
    return full_text


if __name__ == "__main__":
    print("""
╔══════════════════════════════════════════════════════════════════╗
║    Universal Handwriting Recognition - Any File Type             ║
║    Supports: Images, PDF, Word (DOCX), PowerPoint (PPTX)        ║
╚══════════════════════════════════════════════════════════════════╝
    """)
    
    if len(sys.argv) < 2:
        print("Usage: python recognize_any_file.py <file_path> [api_key] [output_path]")
        print("\nExamples:")
        print("   python recognize_any_file.py document.pdf YOUR_API_KEY")
        print("   python recognize_any_file.py notes.docx YOUR_API_KEY output.txt")
        print("   python recognize_any_file.py slides.pptx YOUR_API_KEY")
        print("   python recognize_any_file.py image.jpg YOUR_API_KEY")
        sys.exit(1)
    
    file_path = sys.argv[1]
    api_key = sys.argv[2] if len(sys.argv) > 2 else None
    output_path = sys.argv[3] if len(sys.argv) > 3 else None
    
    recognize_any_file(
        file_path=file_path,
        api_key=api_key,
        output_path=output_path
    )
